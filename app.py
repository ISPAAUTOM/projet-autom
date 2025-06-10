import streamlit as st
import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
import tempfile

# ============================================================================
# CONFIGURATION STREAMLIT
# ============================================================================

st.set_page_config(
    page_title="Modificateur de documents POLYVIA",
    page_icon="📄",
    layout="wide"
)

# ============================================================================
# CONFIGURATION DES STYLES
# ============================================================================

# Logo et positionnement
LOGO_X = Cm(0.65)
LOGO_Y = Cm(1.16)
LOGO_WIDTH = LOGO_HEIGHT = Cm(3.76)

# ============== STYLES POWERPOINT ==============
# 1) Titre
TITRE_POLICE = "Lexend Black"
TITRE_TAILLE = Pt(42)
TITRE_COULEUR = RGBColor(10, 19, 73)  # #0A1349

# 2) Corps
CORPS_POLICE = "Lexend Regular"
CORPS_TAILLE = Pt(22)
CORPS_COULEUR = RGBColor(10, 19, 73)   # #0A1349

# 3) Bullets (≥1)
BULLET_POLICE = "Lexend Light"
BULLET_TAILLE = Pt(18)
BULLET_COULEUR = RGBColor(10, 19, 73) # #0A1349

# 4) Forme décorative (haut-droite)
FORME_DECORATIVE_COULEUR = RGBColor(117, 30, 102)  # #751e66

# Paramètre : ancien logo
MAX_LEFT_LOGO = Cm(2)
MAX_TOP_LOGO = Cm(2)
DOUBLE_LOGO_WIDTH = 2 * LOGO_WIDTH
DOUBLE_LOGO_HEIGHT = 2 * LOGO_HEIGHT

# Limites pour la forme décorative en haut-droite
SEUIL_DROITE = Cm(15)
SEUIL_HAUT = Cm(5)

# ========== STYLES WORD SPÉCIFIQUES ==========

WORD_TITLE_STYLE_NAMES = ["Title", "Titre 1", "Heading 1"]
WORD_SUBTIT_STYLE_NAMES = ["Subtitle", "Titre 2", "Heading 2"]

WORD_TITRE_POLICE = "Lexend Medium"
WORD_TITRE_TAILLE = DocxPt(28)

WORD_SOUS_TITRE_POLICE = "Lexend Medium"
WORD_SOUS_TITRE_TAILLE = DocxPt(14)

WORD_TEXTE_POLICE = "Lexend"
WORD_TEXTE_TAILLE = DocxPt(11)

# ---------------------------------------------------------------------------
# OUTILS PPTX
# ---------------------------------------------------------------------------

def remove_old_logo_if_small_in_corner(shape, progress_text):
    """Si c'est un PICTURE en haut-gauche, <=2x dimension => on supprime."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if (shape.left < MAX_LEFT_LOGO and shape.top < MAX_TOP_LOGO
                and shape.width <= DOUBLE_LOGO_WIDTH
                and shape.height <= DOUBLE_LOGO_HEIGHT):
                shape._element.getparent().remove(shape._element)
                progress_text.text(f"  → Ancien logo supprimé")
                return True
    except Exception as e:
        progress_text.text(f"  ⚠️ Erreur suppression logo: {str(e)}")
    return False

def recolorer_formes_decoratives(objet, progress_text):
    """Parcourt les shapes et sous-shapes, recolore AutoShape en haut-droite."""
    for shape in objet.shapes:
        recolorer_une_forme_recursive(shape, progress_text)

def recolorer_une_forme_recursive(shape, progress_text):
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                recolorer_une_forme_recursive(sub, progress_text)
        else:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and shape.left > SEUIL_DROITE
                and shape.top < SEUIL_HAUT):
                
                # Version corrigée de la recoloration
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = FORME_DECORATIVE_COULEUR
                    progress_text.text(f"  → Forme décorative recolorée")
                except Exception as e:
                    # Si erreur, on essaie une approche alternative
                    try:
                        if hasattr(shape.fill, '_xPr'):
                            shape.fill._xPr.clear()
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = FORME_DECORATIVE_COULEUR
                    except:
                        pass  # On ignore si impossible à recolorer
                        
    except Exception as e:
        # On continue même si une forme pose problème
        pass

def apply_paragraph_style(paragraph, font_name, font_size, font_color):
    try:
        if paragraph.font:
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            if paragraph.font.color:
                paragraph.font.color.rgb = font_color

        for run in paragraph.runs:
            if run.font:
                run.font.name = font_name
                run.font.size = font_size
                if run.font.color:
                    run.font.color.rgb = font_color
    except Exception:
        pass  # On continue si problème de style

def appliquer_style_texte_pptx(text_frame, shape=None, progress_text=None):
    """On désactive l'AutoFit => Forcer 42 pt si titre."""
    if not text_frame:
        return

    try:
        # Désactiver l'AutoFit
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except:
        pass

    force_title = getattr(shape, "_force_title", False) if shape else False
    
    for paragraph in text_frame.paragraphs:
        try:
            if paragraph.level >= 1:
                # bullet => 18
                apply_paragraph_style(paragraph, BULLET_POLICE, BULLET_TAILLE, BULLET_COULEUR)
                if progress_text:
                    progress_text.text(f"    → [Bullet] {paragraph.text[:40]}...")
            else:
                if force_title:
                    apply_paragraph_style(paragraph, TITRE_POLICE, TITRE_TAILLE, TITRE_COULEUR)
                    if progress_text:
                        progress_text.text(f"    → [Titre 42pt] {paragraph.text[:40]}...")
                else:
                    apply_paragraph_style(paragraph, CORPS_POLICE, CORPS_TAILLE, CORPS_COULEUR)
                    if progress_text:
                        progress_text.text(f"    → [Corps 22pt] {paragraph.text[:40]}...")
        except Exception:
            pass  # On continue si problème

def style_table(table):
    try:
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        apply_paragraph_style(paragraph, CORPS_POLICE, CORPS_TAILLE, CORPS_COULEUR)
    except:
        pass

def get_text_content(shape):
    """Retourne le texte d'un shape (strip) ou ""."""
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            return shape.text_frame.text.strip()
    except:
        pass
    return ""

def traiter_pptx(fichier_entree, logo_path, progress_bar, progress_text):
    """Traite un fichier PowerPoint."""
    try:
        progress_text.text("Ouverture du fichier PowerPoint...")
        pres = Presentation(fichier_entree)
        
        total_slides = len(pres.slides)
        master_logo_removed = False

        # MASTER
        progress_text.text("Traitement des masters...")
        for master_idx, master in enumerate(pres.slide_masters, start=1):
            progress_text.text(f"Master {master_idx}")
            recolorer_formes_decoratives(master, progress_text)
            
            any_removed = False
            for shape in list(master.shapes):
                if remove_old_logo_if_small_in_corner(shape, progress_text):
                    any_removed = True

            if any_removed:
                master_logo_removed = True

            # Tri par top + filtre
            text_shapes = [sh for sh in master.shapes if hasattr(sh, "text_frame")]
            text_shapes.sort(key=lambda s: s.top)
            filtered = []
            for sh in text_shapes:
                txt = get_text_content(sh)
                if len(txt) > 3:  # <=3 => ignoré
                    filtered.append(sh)

            if filtered:
                filtered[0]._force_title = True

            for sh in text_shapes:
                appliquer_style_texte_pptx(sh.text_frame, sh, progress_text)

        # SLIDES
        slides_list = list(pres.slides)
        for idx, slide in enumerate(slides_list, start=1):
            progress_bar.progress(idx / total_slides)
            progress_text.text(f"Slide {idx}/{total_slides}")
            
            recolorer_formes_decoratives(slide, progress_text)

            old_removed = False
            for shape in list(slide.shapes):
                if remove_old_logo_if_small_in_corner(shape, progress_text):
                    old_removed = True

            if old_removed or master_logo_removed:
                try:
                    slide.shapes.add_picture(logo_path, LOGO_X, LOGO_Y, 
                                           width=LOGO_WIDTH, height=LOGO_HEIGHT)
                    progress_text.text("  → Nouveau logo inséré")
                except Exception as e:
                    progress_text.text(f"  ⚠️ Erreur insertion logo: {str(e)}")

            text_shapes = []
            for sh in slide.shapes:
                if hasattr(sh, "text_frame"):
                    text_shapes.append(sh)
                if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                    style_table(sh.table)

            text_shapes.sort(key=lambda s: s.top)
            filtered = []
            for sh in text_shapes:
                txt = get_text_content(sh)
                if len(txt) > 3:
                    filtered.append(sh)

            if filtered:
                filtered[0]._force_title = True

            for sh in text_shapes:
                appliquer_style_texte_pptx(sh.text_frame, sh, progress_text)

        # Sauvegarder
        output = io.BytesIO()
        pres.save(output)
        output.seek(0)
        
        progress_bar.progress(1.0)
        progress_text.text("✅ PowerPoint traité avec succès!")
        
        return output

    except Exception as e:
        st.error(f"❌ Erreur PowerPoint: {str(e)}")
        return None

# ============================================================================
# TRAITEMENT WORD
# ============================================================================

def apply_run_style_word(run, font_name, font_size):
    try:
        if run.font:
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = DocxRGBColor(0, 0, 0)
    except:
        pass

def style_word_paragraph_by_name(paragraph):
    try:
        if paragraph.style and paragraph.style.name:
            style_name = paragraph.style.name.lower()
            if any(s.lower() in style_name for s in WORD_TITLE_STYLE_NAMES):
                for run in paragraph.runs:
                    apply_run_style_word(run, WORD_TITRE_POLICE, WORD_TITRE_TAILLE)
                return "TITLE"
            elif any(s.lower() in style_name for s in WORD_SUBTIT_STYLE_NAMES):
                for run in paragraph.runs:
                    apply_run_style_word(run, WORD_SOUS_TITRE_POLICE, WORD_SOUS_TITRE_TAILLE)
                return "SUB"
    except:
        pass
    return None

def appliquer_style_texte_word(paragraph, is_title_fallback=False):
    txt = paragraph.text.strip()
    if not txt:
        return

    style_result = style_word_paragraph_by_name(paragraph)
    if style_result == "TITLE" or style_result == "SUB":
        return

    is_bullet = (txt.startswith("- ") or txt.startswith("* "))
    if is_bullet:
        for run in paragraph.runs:
            apply_run_style_word(run, WORD_SOUS_TITRE_POLICE, WORD_SOUS_TITRE_TAILLE)
    else:
        if is_title_fallback:
            for run in paragraph.runs:
                apply_run_style_word(run, WORD_TITRE_POLICE, WORD_TITRE_TAILLE)
        else:
            for run in paragraph.runs:
                apply_run_style_word(run, WORD_TEXTE_POLICE, WORD_TEXTE_TAILLE)

def traiter_docx(fichier_entree, logo_path, progress_bar, progress_text):
    try:
        progress_text.text("Ouverture du fichier Word...")
        doc = Document(fichier_entree)

        # En-têtes
        progress_text.text("Traitement des en-têtes...")
        for section_idx, section in enumerate(doc.sections, start=1):
            progress_bar.progress(0.3)
            header = section.header
            logo_found = False

            for para in header.paragraphs:
                for run in para.runs:
                    if run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                        logo_found = True
                        run._element.clear()

            if logo_found:
                run = header.paragraphs[0].add_run()
                run.add_picture(logo_path, width=LOGO_WIDTH, height=LOGO_HEIGHT)
                progress_text.text("✅ Nouveau logo ajouté à l'en-tête")

        # 1ère image corps => supprime + insère
        progress_text.text("Recherche et remplacement du logo...")
        progress_bar.progress(0.5)
        found_first_image = False
        for para in doc.paragraphs:
            if not found_first_image:
                for run in para.runs:
                    if run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                        run._element.clear()
                        new_run = para.add_run()
                        new_run.add_picture(logo_path, width=LOGO_WIDTH, height=LOGO_HEIGHT)
                        progress_text.text("✅ Logo remplacé dans le document")
                        found_first_image = True
                        break
            else:
                break

        # Styles
        progress_text.text("Application des styles...")
        progress_bar.progress(0.7)
        found_title = False
        for para in doc.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            if not found_title:
                appliquer_style_texte_word(para, is_title_fallback=True)
                found_title = True
            else:
                appliquer_style_texte_word(para, is_title_fallback=False)

        # Sauvegarder
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        
        progress_bar.progress(1.0)
        progress_text.text("✅ Document Word traité avec succès!")
        
        return output

    except Exception as e:
        st.error(f"❌ Erreur Word: {str(e)}")
        return None

# ============================================================================
# INTERFACE STREAMLIT
# ============================================================================

def main():
    st.title("🎨 Modificateur de documents POLYVIA")
    st.markdown("### Transformez vos présentations PowerPoint et documents Word")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("#### 📁 Document à traiter")
        uploaded_file = st.file_uploader(
            "Glissez votre fichier ici",
            type=['pptx', 'docx'],
            help="Formats supportés : PowerPoint (.pptx) et Word (.docx)"
        )
    
    with col2:
        st.markdown("#### 🖼️ Nouveau logo")
        logo_file = st.file_uploader(
            "Glissez votre logo ici",
            type=['png', 'jpg', 'jpeg'],
            help="Le logo sera redimensionné automatiquement"
        )
    
    if uploaded_file and logo_file:
        st.markdown("---")
        
        # Afficher les infos du fichier
        file_details = {
            "Nom du fichier": uploaded_file.name,
            "Type": uploaded_file.type,
            "Taille": f"{uploaded_file.size / 1024:.1f} KB"
        }
        st.json(file_details)
        
        if st.button("🚀 Lancer le traitement", type="primary"):
            
            # Créer des fichiers temporaires
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_logo:
                tmp_logo.write(logo_file.getbuffer())
                logo_path = tmp_logo.name
            
            # Progress bars
            progress_bar = st.progress(0)
            progress_text = st.empty()
            
            # Traitement selon le type
            if uploaded_file.name.lower().endswith('.pptx'):
                with st.spinner('Traitement du PowerPoint en cours...'):
                    output = traiter_pptx(uploaded_file, logo_path, progress_bar, progress_text)
            else:
                with st.spinner('Traitement du Word en cours...'):
                    output = traiter_docx(uploaded_file, logo_path, progress_bar, progress_text)
            
            # Nettoyer le fichier temporaire
            os.unlink(logo_path)
            
            if output:
                st.success("✅ Traitement terminé avec succès!")
                
                # Bouton de téléchargement
                st.download_button(
                    label="📥 Télécharger le fichier modifié",
                    data=output,
                    file_name=f"modifie_{uploaded_file.name}",
                    mime="application/octet-stream",
                    type="primary"
                )
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #666;'>
        <small>
        💡 Conseils : Le script va automatiquement détecter et remplacer l'ancien logo, 
        appliquer les nouvelles polices et ajuster les tailles de texte.
        </small>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
